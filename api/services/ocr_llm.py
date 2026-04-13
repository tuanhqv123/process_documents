"""OCR service using OpenAI-compatible API (dots.ocr model).

Converts PDF pages to images, sends to vLLM for layout analysis with
bounding boxes, returns structured JSON per page.
"""

import base64
import io
import json
import logging
import os
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Optional

from api.config import settings

logger = logging.getLogger(__name__)

OCR_LAYOUT_PROMPT = """Please output the layout information from the PDF image, including each layout element's bbox, its category, and the corresponding text content within the bbox.

1. Bbox format: [x1, y1, x2, y2]

2. Layout Categories: The possible categories are ['Caption', 'Footnote', 'Formula', 'List-item', 'Page-footer', 'Page-header', 'Picture', 'Section-header', 'Table', 'Text', 'Title'].

3. Text Extraction & Formatting Rules:
    - Picture: For the 'Picture' category, the text field should be omitted.
    - Formula: Format its text as LaTeX.
    - Table: Format its text as HTML.
    - All Others (Text, Title, etc.): Format their text as Markdown.

4. Constraints:
    - The output text must be the original text from the image, with no translation.
    - All layout elements must be sorted according to human reading order.

5. Final Output: The entire output must be a single JSON array."""


def _get_ocr_client():
    """Return an OpenAI client configured from the active OCR key in DB, falling back to env vars."""
    from openai import OpenAI
    try:
        from api.db import get_session, ApiKey
        db = get_session()
        try:
            active = db.query(ApiKey).filter(ApiKey.type == "ocr", ApiKey.is_active == True).first()
            if active and active.base_url:
                return OpenAI(api_key=active.api_key or "0", base_url=active.base_url.rstrip("/"))
        finally:
            db.close()
    except Exception:
        pass
    # Fallback to env config
    if settings.DOTS_OCR_URL:
        return OpenAI(api_key="0", base_url=settings.DOTS_OCR_URL)
    raise ValueError("No active OCR model configured. Go to Settings and set one as active.")


def _get_llm_client():
    """Return an OpenAI client for the active LLM key, or None if none configured."""
    from openai import OpenAI
    try:
        from api.db import get_session, ApiKey
        db = get_session()
        try:
            active = db.query(ApiKey).filter(ApiKey.type == "llm", ApiKey.is_active == True).first()
            if active and active.base_url:
                return OpenAI(api_key=active.api_key or "0", base_url=active.base_url.rstrip("/"))
        finally:
            db.close()
    except Exception:
        pass
    return None


def _get_llm_model_name() -> str:
    """Get model name from active LLM key in DB."""
    try:
        from api.db import get_session, ApiKey
        db = get_session()
        try:
            active = db.query(ApiKey).filter(ApiKey.type == "llm", ApiKey.is_active == True).first()
            if active and active.model_name:
                return active.model_name
        finally:
            db.close()
    except Exception:
        pass
    return getattr(settings, "LLM_MODEL", "gpt-4o")


def _extract_page_text(layout_json: list) -> str:
    """Extract readable text from non-picture blocks for use as page context."""
    parts = []
    for b in layout_json:
        if b.get("category") in ("Picture", "Figure", "Page-header", "Page-footer"):
            continue
        t = (b.get("text") or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts)


def _caption_pictures(layout_json: list, img_b64: str, img_w: int, img_h: int) -> list:
    """Crop each Picture/Figure block and send to the active LLM for captioning.

    Uses surrounding page text as context and responds in the document's language.
    Modifies layout_json in-place: sets block['text'] to the generated caption.
    Returns the (possibly updated) layout_json.
    Silently skips if no LLM is configured or captioning fails.
    """
    from PIL import Image

    picture_indices = [
        i for i, b in enumerate(layout_json)
        if b.get("category") in ("Picture", "Figure") and b.get("bbox")
    ]
    if not picture_indices:
        return layout_json

    llm_client = _get_llm_client()
    if not llm_client:
        logger.info("Caption: no active LLM configured, skipping %d picture(s)", len(picture_indices))
        return layout_json

    logger.info("Caption: captioning %d picture block(s) via LLM", len(picture_indices))

    llm_model = _get_llm_model_name()
    page_context = _extract_page_text(layout_json)

    try:
        img_bytes = base64.b64decode(img_b64)
        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    except Exception as e:
        logger.warning("Caption: could not decode page image: %s", e)
        return layout_json

    context_block = (
        f"Page context (use this to understand the document topic and language):\n{page_context}\n\n"
        if page_context else ""
    )

    caption_prompt = (
        f"{context_block}"
        "Analyze the image above extracted from a document page.\n"
        "1. Identify the image type (chart, diagram, photo, table, illustration, etc.)\n"
        "2. Describe the content in detail\n"
        "3. Explain its meaning in relation to the surrounding document context\n\n"
        "IMPORTANT: Respond in the same language as the page context text above.\n"
        "Format:\nType: [image type]\nDescription: [detailed description and meaning]"
    )

    for i in picture_indices:
        block = layout_json[i]
        x1, y1, x2, y2 = block["bbox"]
        x1, y1 = max(0, int(x1)), max(0, int(y1))
        x2, y2 = min(img_w, int(x2)), min(img_h, int(y2))
        if x2 <= x1 or y2 <= y1:
            continue

        cropped = pil_img.crop((x1, y1, x2, y2))
        buf = io.BytesIO()
        cropped.save(buf, format="PNG")
        crop_b64 = base64.b64encode(buf.getvalue()).decode()

        try:
            resp = llm_client.chat.completions.create(
                model=llm_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{crop_b64}"}},
                            {"type": "text", "text": caption_prompt},
                        ],
                    },
                ],
                max_tokens=512,
                temperature=0.1,
                extra_body={"chat_template_kwargs": {"enable_thinking": False}},
            )
            caption = (resp.choices[0].message.content or "").strip()
            if caption:
                block["text"] = caption
                logger.debug("Captioned picture block %d: %s…", i, caption[:60])
        except Exception as e:
            logger.warning("Caption failed for block %d: %s", i, e)

    return layout_json


def _parse_layout_json(raw_text: str) -> Optional[list]:
    text = raw_text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```\w*\n?", "", text)
        text = re.sub(r"\n?```$", "", text)
        text = text.strip()
    try:
        data = json.loads(text)
        if isinstance(data, list):
            return data
        if isinstance(data, dict) and "layout" in data:
            return data["layout"]
    except json.JSONDecodeError:
        match = re.search(r"\[.*\]", text, re.DOTALL)
        if match:
            try:
                return json.loads(match.group())
            except json.JSONDecodeError:
                pass
    logger.warning("Failed to parse layout JSON (len=%d)", len(raw_text))
    return None


def _clean_markdown(text: str) -> str:
    text = re.sub(r"^(#{1,6})\s+(?:#{1,6}\s+)+", r"\1 ", text, flags=re.MULTILINE)
    text = re.sub(r"^(-\s+){2,}", "- ", text, flags=re.MULTILINE)
    return text


def _layout_to_markdown(layout_json: list) -> str:
    parts = []
    for el in layout_json:
        cat = el.get("category", "Text")
        text = el.get("text", "").strip()
        if not text and cat != "Picture":
            continue
        if cat == "Title":
            parts.append(f"# {text}")
        elif cat == "Section-header":
            parts.append(f"## {text}")
        elif cat == "Formula":
            parts.append(f"$${text}$$")
        elif cat == "Table":
            parts.append(text)
        elif cat == "Picture":
            parts.append(text if text else "[Image]")
        elif cat in ("Page-header", "Page-footer"):
            continue
        elif cat == "Caption":
            parts.append(f"*{text}*")
        elif cat == "Footnote":
            parts.append(f"> {text}")
        elif cat == "List-item":
            parts.append(f"- {text}")
        else:
            parts.append(text)
    return "\n\n".join(parts)


def _is_blank_image(img_b64: str) -> bool:
    """Return True only if the page is truly empty.

    Counts pixels darker than 240 on a 256×256 thumbnail.
    A page needs < 0.3% dark pixels to be considered blank.
    This handles faint text, light logos, lightly printed pages correctly.
    """
    try:
        from PIL import Image
        img_bytes = base64.b64decode(img_b64)
        img = Image.open(io.BytesIO(img_bytes)).convert("L")
        thumb = img.resize((256, 256))
        pixels = thumb.getdata()
        dark = sum(1 for p in pixels if p < 240)
        return dark / len(pixels) < 0.003   # < 0.3% dark pixels = blank
    except Exception:
        return False


def _get_ocr_model_name() -> str:
    """Get model name from active OCR key in DB, fallback to env var."""
    try:
        from api.db import get_session, ApiKey
        db = get_session()
        try:
            active = db.query(ApiKey).filter(ApiKey.type == "ocr", ApiKey.is_active == True).first()
            if active and active.model_name:
                return active.model_name
        finally:
            db.close()
    except Exception:
        pass
    return settings.DOTS_OCR_MODEL


def _ocr_page_vllm(img_b64: str) -> dict:
    """Send a single base64 image to vLLM for layout OCR."""
    last_error = None
    for attempt in range(1 + settings.OCR_MAX_RETRIES):
        try:
            client = _get_ocr_client()
            resp = client.chat.completions.create(
                model=_get_ocr_model_name(),
                messages=[{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                    {"type": "text", "text": OCR_LAYOUT_PROMPT},
                ]}],
                max_completion_tokens=32768,
                temperature=0.1,
            )
            raw = resp.choices[0].message.content or ""
            layout = _parse_layout_json(raw)
            if layout:
                markdown = _clean_markdown(_layout_to_markdown(layout))
                return {"layout_json": layout, "markdown": markdown}
            else:
                return {"layout_json": [], "markdown": _clean_markdown(raw)}
        except Exception as e:
            last_error = e
            if attempt < settings.OCR_MAX_RETRIES:
                logger.warning("vLLM OCR attempt %d failed: %s", attempt + 1, e)
            else:
                logger.error("vLLM OCR failed after %d attempts: %s", attempt + 1, e)
    raise last_error


def pdf_pages_to_b64(file_bytes: bytes, scale: float = 2.0) -> list[tuple]:
    """Convert each PDF page to (base64_png, width, height)."""
    import pypdfium2 as pdfium
    from PIL import Image

    pdf = pdfium.PdfDocument(file_bytes)
    results = []
    for i in range(len(pdf)):
        page = pdf[i]
        bitmap = page.render(scale=scale)
        pil_img = bitmap.to_pil()
        if pil_img.mode != "RGB":
            pil_img = pil_img.convert("RGB")
        w, h = pil_img.size
        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()
        results.append((b64, w, h))
    return results


def pptx_slides_to_b64(file_bytes: bytes) -> list[tuple]:
    """Convert each PPTX slide to (base64_png, width, height).
    Renders slides using python-pptx shapes → PIL image.
    """
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(io.BytesIO(file_bytes))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Render at ~2x scale for good OCR quality
    SCALE = 2.0
    px_w = int(slide_w / Emu(914400) * 96 * SCALE)  # EMU to inches to pixels
    px_h = int(slide_h / Emu(914400) * 96 * SCALE)

    results = []
    for slide in prs.slides:
        img = Image.new("RGB", (px_w, px_h), "white")
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Convert shape position from EMU to pixels
            x = int(shape.left / slide_w * px_w)
            y = int(shape.top / slide_h * px_h)
            w = int(shape.width / slide_w * px_w)

            text = shape.text_frame.text
            if text.strip():
                # Estimate font size from shape height
                font_size = max(12, min(40, int(shape.height / slide_h * px_h * 0.3)))
                try:
                    font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
                except Exception:
                    font = ImageFont.load_default()
                draw.text((x + 10, y + 10), text, fill="black", font=font)

        # Also render images embedded in slides
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    pic = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    x = int(shape.left / slide_w * px_w)
                    y = int(shape.top / slide_h * px_h)
                    w = int(shape.width / slide_w * px_w)
                    h = int(shape.height / slide_h * px_h)
                    pic = pic.resize((w, h))
                    img.paste(pic, (x, y))
                except Exception:
                    pass

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode()
        results.append((b64, px_w, px_h))

    return results


def get_pdf_page_count(file_bytes: bytes) -> int:
    import pypdfium2 as pdfium
    try:
        return len(pdfium.PdfDocument(file_bytes))
    except Exception:
        return 0


def get_pptx_slide_count(file_bytes: bytes) -> int:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        return len(prs.slides)
    except Exception:
        return 0


def ocr_single_page(img_data: tuple, page_num: int) -> dict:
    """OCR one page. img_data = (base64_png, width, height).
    Returns {page, text, page_data, error}.
    Automatically captions Picture/Figure blocks using the active LLM if one is configured.
    """
    try:
        img_b64, img_w, img_h = img_data
        result = _ocr_page_vllm(img_b64)
        layout_json = result.get("layout_json", [])

        # Caption any Picture/Figure blocks via LLM
        layout_json = _caption_pictures(layout_json, img_b64, img_w, img_h)

        # Rebuild markdown after captioning so Picture text is included
        markdown = _clean_markdown(_layout_to_markdown(layout_json))

        page_data = {
            "page": page_num,
            "layout_json": layout_json,
            "markdown": markdown,
            "image_width": img_w,
            "image_height": img_h,
        }
        return {"page": page_num, "text": markdown, "page_data": page_data, "error": None}
    except Exception as e:
        logger.warning("OCR failed page %d: %s", page_num, e)
        return {"page": page_num, "text": "", "page_data": None, "error": str(e)}


def ocr_pdf(
    file_bytes: bytes,
    doc_id: int,
    on_page_done=None,
    cancel_check=None,
    file_type: str = "pdf",
) -> tuple[str, list[dict]]:
    """OCR a full document (PDF or PPTX) using dots.ocr in parallel.

    Saves each page PNG to data/ocr_images/{doc_id}/page_{n}.png.
    Returns (combined_markdown, list_of_page_data).
    """
    images_dir = Path(settings.OCR_IMAGES_DIR) / str(doc_id)
    images_dir.mkdir(parents=True, exist_ok=True)

    try:
        if file_type in ("pptx", "ppt"):
            images_data = pptx_slides_to_b64(file_bytes)
        else:
            images_data = pdf_pages_to_b64(file_bytes)
    except Exception as e:
        logger.error("File to images failed: %s", e)
        raise

    total = len(images_data)
    if total == 0:
        return "", []

    # Save page PNGs to disk before OCR
    for i, (b64, w, h) in enumerate(images_data):
        page_path = images_dir / f"page_{i + 1}.png"
        if not page_path.exists():
            img_bytes = base64.b64decode(b64)
            page_path.write_bytes(img_bytes)

    results = [None] * total
    completed = 0

    with ThreadPoolExecutor(max_workers=settings.OCR_PARALLEL_WORKERS) as executor:
        futures = {}
        for i, img_data in enumerate(images_data):
            if cancel_check and cancel_check():
                break
            img_b64 = img_data[0]
            if _is_blank_image(img_b64):
                results[i] = {"page": i + 1, "text": "", "page_data": None, "error": None}
                completed += 1
                if on_page_done:
                    on_page_done(completed, total)
                continue
            f = executor.submit(ocr_single_page, img_data, i + 1)
            futures[f] = i

        for future in as_completed(futures):
            if cancel_check and cancel_check():
                executor.shutdown(wait=False, cancel_futures=True)
                return "", []

            idx = futures[future]
            try:
                results[idx] = future.result()
            except Exception as e:
                results[idx] = {"page": idx + 1, "text": "", "page_data": None, "error": str(e)}

            completed += 1
            if on_page_done:
                on_page_done(completed, total)

    texts = []
    ocr_pages = []
    for r in results:
        if r and r.get("text"):
            texts.append(f"--- Page {r['page']} ---\n{r['text']}")
        if r and r.get("page_data"):
            ocr_pages.append(r["page_data"])

    return "\n\n".join(texts), ocr_pages


def save_ocr_data(doc_id: int, ocr_pages: list[dict]) -> str:
    """Save OCR page data to JSON file. Returns the file path."""
    data_dir = Path(settings.OCR_DATA_DIR)
    data_dir.mkdir(parents=True, exist_ok=True)
    path = data_dir / f"{doc_id}_ocr.json"
    path.write_text(json.dumps(ocr_pages, ensure_ascii=False), encoding="utf-8")
    return str(path)


def load_ocr_data(doc_id: int) -> list[dict]:
    """Load saved OCR page data from JSON file."""
    path = Path(settings.OCR_DATA_DIR) / f"{doc_id}_ocr.json"
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return []


def update_ocr_block(doc_id: int, page_num: int, block_idx: int, text: str) -> bool:
    """Update the text of a single OCR block in the saved JSON file.
    Returns True if the block was found and updated, False otherwise.
    """
    pages = load_ocr_data(doc_id)
    for page in pages:
        if page.get("page") == page_num:
            layout = page.get("layout_json", [])
            if 0 <= block_idx < len(layout):
                layout[block_idx]["text"] = text
                save_ocr_data(doc_id, pages)
                return True
            return False
    return False


def get_page_image_path(doc_id: int, page_num: int) -> Optional[Path]:
    """Return path to saved page PNG (1-based page_num), or None if not found."""
    path = Path(settings.OCR_IMAGES_DIR) / str(doc_id) / f"page_{page_num}.png"
    return path if path.exists() else None
